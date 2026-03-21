"""
Meeting PPTX Generator — 브랜드 컬러 기반 동적 테마 PPTX 생성.

generate_pptx.js (PptxGenJS)를 python-pptx로 포팅.
슬라이드 타입별 레이아웃을 프로그래밍 방식으로 생성한다.
"""
from __future__ import annotations

import io
import logging
from typing import Any

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

logger = logging.getLogger(__name__)

# ===== 고정 중립색 (브랜드 무관) =====

NEUTRAL = {
    "text_dark": RGBColor(0x1A, 0x1A, 0x1A),
    "text_mid": RGBColor(0x55, 0x55, 0x55),
    "text_light": RGBColor(0x99, 0x99, 0x99),
    "card_bg": RGBColor(0xF8, 0xF8, 0xF8),
    "card_border": RGBColor(0xE5, 0xE5, 0xE5),
    "white": RGBColor(0xFF, 0xFF, 0xFF),
    "dark_bg": RGBColor(0x1A, 0x1A, 0x1A),
}

DEFAULT_BRAND_COLORS = {
    "primary": "#990011",
    "secondary": "#FCF6F5",
    "accent": "#2F3C7E",
}

FONTS = {
    "header": "Arial Black",
    "body": "Arial",
}


# ===== 유틸리티 =====

def _hex_to_rgb(hex_str: str) -> RGBColor:
    """#RRGGBB 또는 RRGGBB → RGBColor."""
    h = hex_str.lstrip("#")
    if len(h) != 6:
        h = "990011"  # fallback
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _lighten(color: RGBColor, factor: float = 0.2) -> RGBColor:
    """색상을 밝게."""
    r = min(255, int(color[0] + (255 - color[0]) * factor))
    g = min(255, int(color[1] + (255 - color[1]) * factor))
    b = min(255, int(color[2] + (255 - color[2]) * factor))
    return RGBColor(r, g, b)


def _tint_bg(color: RGBColor) -> RGBColor:
    """primary에서 아주 연한 배경색 생성."""
    return _lighten(color, 0.92)


def build_theme(brand_colors: dict | None = None) -> dict:
    """brand_colors dict → 완전한 테마 dict."""
    bc = brand_colors or DEFAULT_BRAND_COLORS
    theme = dict(NEUTRAL)
    theme["primary"] = _hex_to_rgb(bc.get("primary", "#990011"))
    theme["secondary"] = _hex_to_rgb(bc.get("secondary", "#FCF6F5"))
    theme["accent"] = _hex_to_rgb(bc.get("accent", "#2F3C7E"))
    theme["primary_light"] = _lighten(theme["primary"], 0.2)
    theme["primary_tint"] = _tint_bg(theme["primary"])
    return theme


def _add_rect(slide, left, top, width, height, fill_color=None, line_color=None, line_width=None):
    """사각형 shape 추가."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height),
    )
    shape.shadow.inherit = False
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color and line_width:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape


def _add_text(slide, left, top, width, height, text, font_size=12,
              font_name=None, color=None, bold=False, align=PP_ALIGN.LEFT,
              valign=MSO_ANCHOR.TOP, char_spacing=None, line_spacing=None):
    """텍스트박스 추가."""
    txbox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height),
    )
    tf = txbox.text_frame
    tf.word_wrap = True
    tf.auto_size = None

    p = tf.paragraphs[0]
    p.text = str(text)
    p.font.size = Pt(font_size)
    p.font.name = font_name or FONTS["body"]
    if color:
        p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align

    tf.paragraphs[0].space_before = Pt(0)
    tf.paragraphs[0].space_after = Pt(0)

    if char_spacing is not None:
        p.font._element.attrib['{http://schemas.openxmlformats.org/drawingml/2006/main}spc'] = str(int(char_spacing * 100))

    if line_spacing:
        p.line_spacing = Pt(line_spacing)

    # vertical alignment
    from pptx.oxml.ns import qn
    bodyPr = tf._txBody.find(qn('a:bodyPr'))
    if bodyPr is not None:
        anchor_map = {
            MSO_ANCHOR.TOP: 't',
            MSO_ANCHOR.MIDDLE: 'ctr',
            MSO_ANCHOR.BOTTOM: 'b',
        }
        bodyPr.set('anchor', anchor_map.get(valign, 't'))

    return txbox


def _add_multiline_text(slide, left, top, width, height, lines, font_size=12,
                        font_name=None, color=None, bold=False, align=PP_ALIGN.LEFT,
                        line_spacing_multiple=1.3):
    """여러 줄 텍스트박스."""
    txbox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height),
    )
    tf = txbox.text_frame
    tf.word_wrap = True
    tf.auto_size = None

    for i, line_data in enumerate(lines):
        if isinstance(line_data, str):
            line_text, line_color, line_bold = line_data, color, bold
        else:
            line_text = line_data.get("text", "")
            line_color = line_data.get("color", color)
            line_bold = line_data.get("bold", bold)

        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        p.text = str(line_text)
        p.font.size = Pt(font_size)
        p.font.name = font_name or FONTS["body"]
        if line_color:
            p.font.color.rgb = line_color
        p.font.bold = line_bold
        p.alignment = align

    return txbox


def _add_card(slide, x, y, w, h, theme):
    """카드 배경 (회색 배경 + 테두리)."""
    _add_rect(slide, x, y, w, h, fill_color=theme["card_bg"],
              line_color=theme["card_border"], line_width=0.5)


def _add_accent_bar(slide, x, y, w, h, color):
    """색상 바."""
    _add_rect(slide, x, y, w, h, fill_color=color)


def _add_line(slide, x, y, w, color, width_pt=0.75):
    """수평선."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Pt(width_pt),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


# ===== 슬라이드 헤더 (공통) =====

def _add_slide_header(slide, data: dict, theme: dict):
    """콘텐츠 슬라이드 공통 헤더: 좌측 빨간 바 + 섹션명 + 페이지 번호 + 구분선."""
    # 좌측 accent bar
    _add_accent_bar(slide, 0.45, 0.35, 0.06, 0.4, theme["primary"])

    # 섹션 타이틀
    _add_text(slide, 0.6, 0.3, 7, 0.5,
              data.get("section_title", ""),
              font_size=22, font_name=FONTS["header"],
              color=theme["text_dark"], bold=True)

    # 페이지 번호
    _add_text(slide, 8.5, 0.3, 1.2, 0.4,
              data.get("page_label", ""),
              font_size=10, color=theme["text_light"], align=PP_ALIGN.RIGHT)

    # 구분선
    _add_line(slide, 0.45, 0.8, 9.1, theme["card_border"])


# ===== 슬라이드 생성 함수 =====

def _create_cover_slide(prs, data: dict, theme: dict):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 우측 상단 장식 바
    _add_accent_bar(slide, 9.7, 0, 0.3, 0.8, theme["primary"])

    # 라벨
    _add_text(slide, 0.8, 1.2, 8, 0.35,
              data.get("label", "ADVERTISING STRATEGY PROPOSAL"),
              font_size=13, color=theme["primary"], bold=True, char_spacing=3)

    # 메인 타이틀
    _add_text(slide, 0.8, 1.7, 8, 0.9,
              data.get("title", ""),
              font_size=44, font_name=FONTS["header"],
              color=theme["text_dark"], bold=True)

    # 서브 타이틀 (primary 색상)
    _add_text(slide, 0.8, 2.6, 8, 0.8,
              data.get("title_accent", "광고전략 제안서"),
              font_size=42, font_name=FONTS["header"],
              color=theme["primary"], bold=True)

    # 부제
    if data.get("subtitle_line1"):
        _add_multiline_text(
            slide, 0.8, 3.5, 8, 0.7,
            [data.get("subtitle_line1", ""), data.get("subtitle_line2", "")],
            font_size=15, color=theme["text_mid"],
            line_spacing_multiple=1.4,
        )

    # 하단 구분선
    _add_line(slide, 0.8, 4.3, 8.4, theme["card_border"], 1)

    # 메타 정보
    meta = data.get("_meta", {})
    meta_items = [
        ("CLIENT", meta.get("client", "")),
        ("DATE", meta.get("project_date", "")),
        ("PREPARED BY", meta.get("prepared_by", "")),
    ]
    for i, (label, value) in enumerate(meta_items):
        mx = 0.8 + i * 3
        _add_text(slide, mx, 4.5, 2.5, 0.25,
                  label, font_size=8, color=theme["text_light"],
                  bold=True, char_spacing=2)
        _add_text(slide, mx, 4.75, 2.5, 0.3,
                  value, font_size=11, color=theme["text_dark"])

    # 저작권
    _add_text(slide, 0.8, 5.2, 8, 0.25,
              "CONFIDENTIAL & PROPRIETARY",
              font_size=7, color=theme["text_light"])


def _create_pillars_slide(prs, data: dict, theme: dict):
    """pillars_3col, insight_cards, challenge_solution, promise_rtb, creative_grid 공통."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_slide_header(slide, data, theme)

    # 코어 블록 (회색 배경)
    _add_rect(slide, 0.45, 1.0, 9.1, 1.7, fill_color=theme["card_bg"])

    # 코어 라벨
    _add_text(slide, 0.7, 1.05, 4, 0.22,
              data.get("core_label", "CORE OBJECTIVE"),
              font_size=9, color=theme["primary"], bold=True, char_spacing=2)

    # 헤드라인
    _add_text(slide, 0.7, 1.3, 8.5, 0.7,
              data.get("headline", ""),
              font_size=20, font_name=FONTS["header"],
              color=theme["text_dark"], bold=True)

    # 서브 헤드라인
    _add_text(slide, 0.7, 2.05, 8.5, 0.45,
              data.get("sub_headline", ""),
              font_size=11, color=theme["text_mid"])

    # pillars 라벨
    if data.get("pillars_label"):
        _add_text(slide, 0.45, 2.7, 4, 0.3,
                  data["pillars_label"],
                  font_size=11, color=theme["text_dark"], bold=True)

    # 3열 카드
    pillars = data.get("pillars", [])
    card_w = 2.85
    card_h = 1.75
    card_y = 3.05
    gap = 0.2

    for i, pillar in enumerate(pillars[:3]):
        card_x = 0.45 + i * (card_w + gap)

        # 카드 배경
        _add_card(slide, card_x, card_y, card_w, card_h, theme)

        # 상단 컬러 바
        bar_color = theme["primary"] if i < 2 else theme["text_dark"]
        _add_accent_bar(slide, card_x, card_y, card_w, 0.05, bar_color)

        # 번호 + 제목
        num = pillar.get("number", f"0{i + 1}")
        title = pillar.get("title_ko", "")
        _add_text(slide, card_x + 0.15, card_y + 0.2, card_w - 0.3, 0.3,
                  f"{num}. {title}",
                  font_size=13, color=theme["text_dark"], bold=True)

        # 태그라인
        tag = pillar.get("tag_line", "")
        if tag:
            _add_rect(slide, card_x + 0.15, card_y + 0.6, card_w - 0.3, 0.3,
                      fill_color=theme["card_bg"],
                      line_color=theme["card_border"], line_width=0.5)
            _add_text(slide, card_x + 0.15, card_y + 0.6, card_w - 0.3, 0.3,
                      tag, font_size=9, color=theme["primary"],
                      bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

        # 설명
        _add_text(slide, card_x + 0.15, card_y + 1.05, card_w - 0.3, 0.75,
                  pillar.get("description", ""),
                  font_size=9.5, color=theme["text_mid"])

    # 하단 메트릭
    metrics = data.get("bottom_metrics", [])
    if metrics:
        _add_text(slide, 0.45, 5.15, 1.2, 0.25,
                  "KEY METRICS", font_size=8, color=theme["text_light"], bold=True)
        for i, metric in enumerate(metrics[:3]):
            dot_color = theme["primary"] if i == 0 else theme["text_dark"]
            _add_text(slide, 1.7 + i * 2.3, 5.15, 2.2, 0.25,
                      f"●  {metric}", font_size=8, color=dot_color)

    # 하단 인용구
    quote = data.get("bottom_quote", "")
    if quote:
        _add_text(slide, 0.45, 5.1, 9.1, 0.35,
                  f"\u275d  {quote}  \u275e",
                  font_size=9.5, color=theme["primary"],
                  bold=True, align=PP_ALIGN.CENTER)


def _create_comparison_slide(prs, data: dict, theme: dict):
    """comparison_lr — 좌우 대비 레이아웃."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_slide_header(slide, data, theme)

    # 비교 헤드라인
    if data.get("comparison_title"):
        _add_text(slide, 0.45, 0.9, 9.1, 0.5,
                  data["comparison_title"],
                  font_size=18, font_name=FONTS["header"],
                  color=theme["text_dark"], bold=True)

    left = data.get("left", {})
    right = data.get("right", {})
    col_w = 4.35
    col_y = 1.5
    col_h = 3.5

    # 좌측 패널
    _add_card(slide, 0.45, col_y, col_w, col_h, theme)
    _add_text(slide, 0.65, col_y + 0.15, 4, 0.25,
              left.get("label", "MARKET REALITY"),
              font_size=9, color=theme["text_mid"], bold=True, char_spacing=2)
    _add_text(slide, 0.65, col_y + 0.4, col_w - 0.4, 0.4,
              left.get("headline", ""),
              font_size=16, font_name=FONTS["header"],
              color=theme["text_dark"], bold=True)

    for i, item in enumerate(left.get("items", [])[:3]):
        item_y = col_y + 1.0 + i * 0.7
        _add_text(slide, 0.65, item_y, col_w - 0.4, 0.25,
                  item.get("title", ""),
                  font_size=11, color=theme["text_dark"], bold=True)
        _add_text(slide, 0.65, item_y + 0.25, col_w - 0.4, 0.35,
                  item.get("description", ""),
                  font_size=9, color=theme["text_mid"])

    if left.get("bottom_box"):
        _add_rect(slide, 0.55, col_y + col_h - 0.5, col_w - 0.2, 0.4,
                  fill_color=theme["card_bg"])
        _add_text(slide, 0.65, col_y + col_h - 0.5, col_w - 0.4, 0.4,
                  left["bottom_box"],
                  font_size=9, color=theme["text_mid"],
                  align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

    # 화살표
    _add_text(slide, 4.6, 3.0, 0.5, 0.5,
              "\u25b6", font_size=18, color=theme["primary"],
              align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

    # 우측 패널 (primary 테두리)
    _add_rect(slide, 5.2, col_y, col_w, col_h,
              fill_color=theme["primary_tint"],
              line_color=theme["primary"], line_width=1)

    _add_text(slide, 5.4, col_y + 0.15, 4, 0.25,
              right.get("label", "STRATEGIC OPPORTUNITY"),
              font_size=9, color=theme["primary"], bold=True, char_spacing=2)
    _add_text(slide, 5.4, col_y + 0.4, col_w - 0.4, 0.4,
              right.get("headline", ""),
              font_size=16, font_name=FONTS["header"],
              color=theme["text_dark"], bold=True)

    for i, item in enumerate(right.get("items", [])[:3]):
        item_y = col_y + 1.0 + i * 0.7
        _add_text(slide, 5.4, item_y, col_w - 0.4, 0.25,
                  item.get("title", ""),
                  font_size=11, color=theme["text_dark"], bold=True)
        _add_text(slide, 5.4, item_y + 0.25, col_w - 0.4, 0.35,
                  item.get("description", ""),
                  font_size=9, color=theme["text_mid"])

    if right.get("bottom_box"):
        _add_rect(slide, 5.3, col_y + col_h - 0.5, col_w - 0.2, 0.4,
                  fill_color=theme["primary"])
        _add_text(slide, 5.4, col_y + col_h - 0.5, col_w - 0.4, 0.4,
                  right["bottom_box"],
                  font_size=9, color=theme["white"],
                  align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE,
                  bold=True)


def _create_process_slide(prs, data: dict, theme: dict):
    """process_steps — 3단계 프로세스 카드."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_slide_header(slide, data, theme)

    # 전략명
    strategy = data.get("strategy_name", "")
    if strategy:
        parts = strategy.split(":", 1)
        if len(parts) == 2:
            # 첫 부분은 primary, 나머지는 text_dark
            _add_text(slide, 0.45, 0.9, 9, 0.45,
                      strategy,
                      font_size=18, font_name=FONTS["header"],
                      color=theme["primary"], bold=True)
        else:
            _add_text(slide, 0.45, 0.9, 9, 0.45,
                      strategy,
                      font_size=18, font_name=FONTS["header"],
                      color=theme["text_dark"], bold=True)

    if data.get("strategy_subtitle"):
        _add_text(slide, 0.45, 1.35, 9, 0.3,
                  data["strategy_subtitle"],
                  font_size=11, color=theme["text_mid"])

    # 3단계 카드
    steps = data.get("steps", [])
    step_w = 2.85
    step_h = 2.8
    step_y = 1.85
    gap = 0.2

    for i, step in enumerate(steps[:3]):
        sx = 0.45 + i * (step_w + gap)
        is_highlight = (i == 1)

        # 카드 배경
        if is_highlight:
            _add_rect(slide, sx, step_y, step_w, step_h,
                      fill_color=theme["white"],
                      line_color=theme["primary"], line_width=2)
        else:
            _add_card(slide, sx, step_y, step_w, step_h, theme)

        # 상단 바
        bar_colors = [theme["card_bg"], theme["primary"], theme["text_dark"]]
        _add_accent_bar(slide, sx, step_y, step_w, 0.04, bar_colors[min(i, 2)])

        # STEP 라벨
        _add_text(slide, sx + 0.15, step_y + 0.15, 1.5, 0.2,
                  step.get("step_number", f"STEP 0{i + 1}"),
                  font_size=8, color=theme["primary"] if is_highlight else theme["text_light"],
                  bold=True, char_spacing=2)

        # 영문명
        _add_text(slide, sx + 0.15, step_y + 0.35, 1.8, 0.35,
                  step.get("name_en", ""),
                  font_size=18, font_name=FONTS["header"],
                  color=theme["text_dark"], bold=True)

        # 한글명
        _add_text(slide, sx + 1.85, step_y + 0.4, 0.85, 0.25,
                  step.get("name_ko", ""),
                  font_size=8.5, color=theme["text_light"], align=PP_ALIGN.RIGHT)

        # 인용구
        quote = step.get("quote", "")
        if quote:
            bg = theme["primary_tint"] if is_highlight else theme["card_bg"]
            _add_rect(slide, sx + 0.15, step_y + 0.85, step_w - 0.3, 0.5,
                      fill_color=bg)
            _add_text(slide, sx + 0.25, step_y + 0.85, step_w - 0.5, 0.5,
                      f'"{quote}"',
                      font_size=10, color=theme["primary"] if is_highlight else theme["text_dark"],
                      bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

        # 설명
        _add_text(slide, sx + 0.15, step_y + 1.45, step_w - 0.3, 0.5,
                  step.get("description", ""),
                  font_size=9, color=theme["text_mid"])

        # 구분선
        _add_line(slide, sx + 0.15, step_y + 2.0, step_w - 0.3, theme["card_border"], 0.5)

        # 실행 항목
        for j, ex in enumerate(step.get("executions", [])[:3]):
            _add_text(slide, sx + 0.15, step_y + 2.1 + j * 0.22, step_w - 0.3, 0.2,
                      f"● {ex.get('label', '')}: {ex.get('content', '')}",
                      font_size=8, color=theme["text_mid"])

        # 화살표
        if i < len(steps) - 1 and i < 2:
            _add_text(slide, sx + step_w + 0.02, step_y + 1.2, 0.16, 0.3,
                      "\u25b6", font_size=12, color=theme["text_light"],
                      align=PP_ALIGN.CENTER)

    # 하단 인용구
    if data.get("bottom_quote"):
        _add_text(slide, 0.45, 5.05, 9.1, 0.35,
                  f"\u275d  {data['bottom_quote']}  \u275e",
                  font_size=9.5, color=theme["primary"],
                  bold=True, align=PP_ALIGN.CENTER)


def _create_closing_slide(prs, data: dict, theme: dict):
    """closing — 마무리 슬라이드."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 헤더
    _add_accent_bar(slide, 0.45, 0.35, 0.06, 0.4, theme["primary"])
    _add_text(slide, 0.6, 0.3, 4, 0.5,
              "Closing", font_size=22, font_name=FONTS["header"],
              color=theme["text_dark"], bold=True)
    _add_text(slide, 8.5, 0.3, 1.2, 0.4,
              data.get("page_label", "10 / 10"),
              font_size=10, color=theme["text_light"], align=PP_ALIGN.RIGHT)
    _add_line(slide, 0.45, 0.8, 9.1, theme["card_border"])

    # 메인 메시지
    headline = data.get("closing_headline", "")
    accent = data.get("closing_accent", "")

    _add_multiline_text(
        slide, 1, 1.5, 8, 1,
        [
            {"text": f"{headline} ", "color": theme["text_dark"], "bold": True},
            {"text": accent, "color": theme["primary"], "bold": True},
        ],
        font_size=36, font_name=FONTS["header"], align=PP_ALIGN.CENTER,
    )

    # 감사 메시지
    _add_text(slide, 1, 2.5, 8, 0.5,
              data.get("sub_message", "경청해 주셔서 감사합니다."),
              font_size=18, color=theme["text_mid"], align=PP_ALIGN.CENTER)

    # Next Steps
    steps = data.get("next_steps", [])
    if steps:
        _add_line(slide, 1.5, 3.2, 7, theme["primary"], 2)
        _add_rect(slide, 1.5, 3.2, 7, 1.5, fill_color=theme["card_bg"])
        _add_text(slide, 3.5, 3.3, 3, 0.3,
                  "NEXT STEPS", font_size=11, color=theme["primary"],
                  bold=True, align=PP_ALIGN.CENTER, char_spacing=3)

        for i, step in enumerate(steps[:3]):
            sx = 1.7 + i * 2.3
            _add_text(slide, sx, 3.7, 2, 0.25,
                      step.get("title", f"Step 0{i + 1}"),
                      font_size=11, color=theme["text_dark"], bold=True)
            _add_text(slide, sx, 3.95, 2, 0.25,
                      step.get("description", ""),
                      font_size=9, color=theme["text_mid"])
            _add_text(slide, sx, 4.2, 2, 0.2,
                      step.get("timeline", ""),
                      font_size=8, color=theme["text_light"])

            if i < len(steps) - 1 and i < 2:
                _add_text(slide, sx + 1.95, 3.85, 0.3, 0.3,
                          "\u25b6", font_size=10, color=theme["text_light"],
                          align=PP_ALIGN.CENTER)

    # 연락처
    contact = data.get("contact", {})
    if contact.get("email") or contact.get("phone"):
        _add_rect(slide, 2.5, 4.85, 5, 0.45,
                  fill_color=theme["white"],
                  line_color=theme["card_border"], line_width=0.5)
        email = contact.get("email", "")
        phone = contact.get("phone", "")
        _add_text(slide, 2.5, 4.85, 5, 0.45,
                  f"\u2709  {email}    \u260e  {phone}",
                  font_size=10, color=theme["text_mid"],
                  align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)


# ===== 메인 함수 =====

# 타입 → 생성 함수 dispatch
_SLIDE_CREATORS = {
    "cover": _create_cover_slide,
    "comparison_lr": _create_comparison_slide,
    "process_steps": _create_process_slide,
    "closing": _create_closing_slide,
    # 아래 타입들은 모두 pillars 레이아웃으로 처리
    "pillars_3col": _create_pillars_slide,
    "insight_cards": _create_pillars_slide,
    "challenge_solution": _create_pillars_slide,
    "target_segments": _create_pillars_slide,
    "promise_rtb": _create_pillars_slide,
    "creative_grid": _create_pillars_slide,
}


def generate_meeting_pptx(slides_data: dict) -> bytes:
    """
    구조화된 slides.json → PPTX bytes (in-memory).

    Args:
        slides_data: {"meta": {..., "brand_colors": {...}}, "slides": [...]}

    Returns:
        PPTX 파일의 bytes
    """
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    meta = slides_data.get("meta", {})
    brand_colors = meta.get("brand_colors")
    theme = build_theme(brand_colors)

    prs.core_properties.author = meta.get("prepared_by", "FLUX System")
    prs.core_properties.title = f"{meta.get('client', '')} 광고전략 제안서"

    slides = slides_data.get("slides", [])
    for slide_data in slides:
        # meta를 각 슬라이드에 전달
        slide_data["_meta"] = meta

        slide_type = slide_data.get("type", "pillars_3col")
        creator = _SLIDE_CREATORS.get(slide_type, _create_pillars_slide)

        try:
            creator(prs, slide_data, theme)
        except Exception:
            logger.exception("Error creating slide %s (type=%s), skipping",
                             slide_data.get("slide_number"), slide_type)
            continue

    # bytes로 출력
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
