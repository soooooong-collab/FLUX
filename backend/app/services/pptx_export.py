"""
PPTX Export — Template-based PowerPoint generation.

Workflow:
  1. Unpack template PPTX into temp directory
  2. For each slide in the content JSON:
     - Pick the best template slide for the slide_pattern
     - Duplicate it via add_slide
     - Inject text content via slide_injector
  3. Remove original template slides (keep only generated ones)
  4. Clean orphaned resources
  5. Pack back into PPTX
  6. Return bytes
"""
from __future__ import annotations

import io
import logging
import re
import shutil
import tempfile
from pathlib import Path
from typing import Any

from scripts.office.unpack import unpack
from scripts.office.pack import pack
from scripts.add_slide import add_slide, _next_rid, _next_slide_id
from scripts.clean import clean
from app.services.slide_injector import inject_content

logger = logging.getLogger(__name__)

# ──────────────────────────────────────────────
# Template configuration
# ──────────────────────────────────────────────

# Default template path (relative to backend/)
DEFAULT_TEMPLATE = Path(__file__).resolve().parents[2] / "data" / "[한컴] Life Plus_최종제출_1117.pptx"

# Map slide_pattern → template slide number to clone
PATTERN_SLIDE_MAP: dict[str, int] = {
    "cover": 1,          # Large centered text
    "statement": 5,      # Impact statement with highlight
    "title_body": 4,     # Title + body text
    "quote": 28,         # Image bg + quote text overlay
    "diagram": 31,       # Multi-shape diagram
    "comparison": 17,    # Comparison / before-after
    "narrative": 33,     # Multi-line poetic text
    "reveal": 34,        # Large concept word (English)
    "reveal_kr": 35,     # Large concept word (Korean) + decorative quotes
}

# Default pattern when slide_pattern is missing or unknown
DEFAULT_PATTERN = "statement"


def _read_xml(path: Path) -> str:
    return path.read_text(encoding="utf-8")


def _write_xml(path: Path, content: str) -> None:
    path.write_text(content, encoding="utf-8")


def _update_presentation_xml(
    unpacked: Path,
    new_slides: list[tuple[str, str]],
    original_slide_count: int,
) -> None:
    """Update presentation.xml to:
    1. Remove original template slides from sldIdLst
    2. Add new slides to sldIdLst in order

    Args:
        new_slides: list of (slide_filename, rId) tuples in desired order
        original_slide_count: number of slides in the original template
    """
    pres_path = unpacked / "ppt" / "presentation.xml"
    if not pres_path.exists():
        return

    pres = _read_xml(pres_path)

    # Find sldIdLst and clear it
    sld_id_lst_match = re.search(
        r"(<p:sldIdLst[^>]*>)(.*?)(</p:sldIdLst>)", pres, re.DOTALL
    )
    if not sld_id_lst_match:
        return

    # Build new sldIdLst entries
    pres_rels_path = unpacked / "ppt" / "_rels" / "presentation.xml.rels"
    pres_rels = _read_xml(pres_rels_path) if pres_rels_path.exists() else ""

    new_entries = []
    base_id = 256  # Standard starting ID for slides

    for i, (slide_name, _) in enumerate(new_slides):
        slide_id = base_id + i

        # Find the rId for this slide in presentation.xml.rels
        rid_match = re.search(
            rf'<Relationship[^>]*Id="(rId\d+)"[^>]*Target="slides/{re.escape(slide_name)}"',
            pres_rels,
        )
        if rid_match:
            rid = rid_match.group(1)
        else:
            continue

        new_entries.append(f'<p:sldId id="{slide_id}" r:id="{rid}"/>')

    # Replace sldIdLst content
    new_sld_id_lst = (
        sld_id_lst_match.group(1)
        + "\n".join(new_entries)
        + sld_id_lst_match.group(3)
    )
    pres = pres[: sld_id_lst_match.start()] + new_sld_id_lst + pres[sld_id_lst_match.end() :]

    _write_xml(pres_path, pres)


def _set_slide_size_16x9(unpacked: Path) -> None:
    """Change slide size from 4:3 to 16:9 (12192000 x 6858000 EMU = 13.333\" x 7.5\")."""
    pres_path = unpacked / "ppt" / "presentation.xml"
    if not pres_path.exists():
        return

    pres = _read_xml(pres_path)
    # Replace sldSz: 4:3 (9144000x6858000) → 16:9 (12192000x6858000)
    pres = re.sub(
        r'<p:sldSz[^/]*/>' ,
        '<p:sldSz cx="12192000" cy="6858000" type="custom"/>',
        pres,
    )
    _write_xml(pres_path, pres)
    logger.info("Slide size set to 16:9 (12192000 x 6858000 EMU)")


def _remove_original_slides(unpacked: Path, original_slides: set[str]) -> None:
    """Remove original template slides from the unpacked directory."""
    slides_dir = unpacked / "ppt" / "slides"
    rels_dir = slides_dir / "_rels"

    for slide_name in original_slides:
        slide_path = slides_dir / slide_name
        if slide_path.exists():
            slide_path.unlink()

        rels_path = rels_dir / f"{slide_name}.rels"
        if rels_path.exists():
            rels_path.unlink()

    # Remove rels entries for original slides from presentation.xml.rels
    pres_rels_path = unpacked / "ppt" / "_rels" / "presentation.xml.rels"
    if pres_rels_path.exists():
        pres_rels = _read_xml(pres_rels_path)
        for slide_name in original_slides:
            pres_rels = re.sub(
                rf'<Relationship[^>]*Target="slides/{re.escape(slide_name)}"[^/]*/>\s*',
                "",
                pres_rels,
            )
        _write_xml(pres_rels_path, pres_rels)

    # Remove Content_Types entries for original slides
    ct_path = unpacked / "[Content_Types].xml"
    if ct_path.exists():
        ct = _read_xml(ct_path)
        for slide_name in original_slides:
            ct = re.sub(
                rf'<Override[^>]*PartName="/ppt/slides/{re.escape(slide_name)}"[^/]*/>\s*',
                "",
                ct,
            )
        _write_xml(ct_path, ct)


def generate_pptx(
    brand_name: str,
    slides: list[dict],
    template_path: str | Path | None = None,
) -> bytes:
    """Generate a PPTX file from slide data using the template-based approach.

    Args:
        brand_name: Brand name (used for fallback text)
        slides: List of slide dicts with slide_pattern, title, body, etc.
        template_path: Optional custom template path. Defaults to bundled template.

    Returns:
        PPTX file as bytes
    """
    template = Path(template_path) if template_path else DEFAULT_TEMPLATE

    if not template.exists():
        logger.error("Template not found: %s, falling back to programmatic generation", template)
        return _fallback_generate(brand_name, slides)

    # Create temp directory for work
    with tempfile.TemporaryDirectory(prefix="flux_pptx_") as tmpdir:
        tmpdir = Path(tmpdir)
        unpacked_dir = tmpdir / "unpacked"
        output_path = tmpdir / "output.pptx"

        # Step 1: Unpack template
        unpack(str(template), str(unpacked_dir))

        # Step 1.5: Set slide size to 16:9
        _set_slide_size_16x9(unpacked_dir)

        slides_dir = unpacked_dir / "ppt" / "slides"

        # Record original template slides
        original_slides = {
            f.name
            for f in slides_dir.glob("slide*.xml")
            if re.match(r"slide\d+\.xml", f.name)
        }
        original_count = len(original_slides)

        # Step 2: For each content slide, duplicate the appropriate template slide
        new_slide_info: list[tuple[str, str]] = []  # (new_slide_name, pattern)

        for slide_data in slides:
            pattern = slide_data.get("slide_pattern", DEFAULT_PATTERN)

            # Also support legacy "layout" field
            if pattern == DEFAULT_PATTERN and "layout" in slide_data:
                pattern = _map_layout_to_pattern(slide_data["layout"])

            template_slide_num = PATTERN_SLIDE_MAP.get(pattern)
            if template_slide_num is None:
                template_slide_num = PATTERN_SLIDE_MAP[DEFAULT_PATTERN]
                logger.warning(
                    "Unknown pattern '%s', using default (statement/slide%d)",
                    pattern, template_slide_num,
                )

            source_slide = f"slide{template_slide_num}.xml"

            # Check source exists
            if not (slides_dir / source_slide).exists():
                logger.warning("Source slide %s not found, skipping", source_slide)
                continue

            # Duplicate the template slide
            new_slide_name = add_slide(str(unpacked_dir), source_slide)
            new_slide_path = slides_dir / new_slide_name

            # Step 3: Inject content
            inject_content(new_slide_path, pattern, slide_data)

            new_slide_info.append((new_slide_name, pattern))

        if not new_slide_info:
            logger.error("No slides generated, returning fallback")
            return _fallback_generate(brand_name, slides)

        # Step 4: Update presentation.xml to show only new slides in order
        _update_presentation_xml(unpacked_dir, new_slide_info, original_count)

        # Step 5: Remove original template slides
        _remove_original_slides(unpacked_dir, original_slides)

        # Step 6: Clean orphaned resources
        clean(str(unpacked_dir))

        # Step 7: Pack into PPTX
        pack(str(unpacked_dir), str(output_path))

        # Return bytes
        return output_path.read_bytes()


def _map_layout_to_pattern(layout: str) -> str:
    """Map legacy layout names to slide_pattern names."""
    mapping = {
        "title_content": "title_body",
        "section_divider": "statement",
        "storytelling_single": "narrative",
        "two_column_compare": "comparison",
        "diagram_flow": "diagram",
        "concept_reveal": "reveal",
        "golden_circle": "diagram",
        "data_table": "title_body",
        "quote_highlight": "quote",
    }
    return mapping.get(layout, DEFAULT_PATTERN)


def _fallback_generate(brand_name: str, slides: list[dict]) -> bytes:
    """Minimal fallback using python-pptx if template is unavailable."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        for slide_data in slides:
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            fill = slide.background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(0x1A, 0x1A, 0x2E)

            title = slide_data.get("title", "")
            body = slide_data.get("body", "")

            if title:
                box = slide.shapes.add_textbox(Inches(0.5), Inches(1), Inches(12), Inches(1))
                tf = box.text_frame
                tf.paragraphs[0].text = title
                tf.paragraphs[0].font.size = Pt(32)
                tf.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                tf.paragraphs[0].font.bold = True

            if body:
                box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12), Inches(4))
                tf = box.text_frame
                tf.word_wrap = True
                tf.paragraphs[0].text = body[:2000]
                tf.paragraphs[0].font.size = Pt(14)
                tf.paragraphs[0].font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)

        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)
        return buf.read()

    except ImportError:
        logger.error("python-pptx not installed, cannot generate fallback PPTX")
        return b""
